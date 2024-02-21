import tkinter as tk
from tkinter import ttk
from tkinter import (
    filedialog, messagebox, Toplevel, Entry, Label, Button, Checkbutton, IntVar, simpledialog, Listbox
)
from tkinter.ttk import Separator, Style, Progressbar
import threading
import os
import json
import re
import flask
from flask_httpauth import HTTPBasicAuth
from threading import Thread
import concurrent.futures
from cryptography.fernet import Fernet
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import docx
import openpyxl
from pptx import Presentation
import markdown2
from bs4 import BeautifulSoup
import pythoncom
from loguru import logger
import logging

# Logger configuration
logger.add("pii_scan_log_{time}.txt", rotation="1 week", compression="zip")

# Flask app and basic auth setup
app = flask.Flask(__name__)
auth = HTTPBasicAuth()

# Environment variables for admin credentials
admin_username = os.environ.get('ADMIN_USERNAME', 'default_admin')
admin_password = os.environ.get('ADMIN_PASSWORD', 'default_secret')
users = {admin_username: admin_password}

# Flask app configuration
app.config['SECRET_KEY'] = os.environ.get('FLASK_SECRET_KEY', 'default_secret_key')

# Basic authentication verification
@auth.verify_password
def verify_password(username, password):
    return username in users and users[username] == password

# Encryption module class definition
class EncryptionModule:
    def __init__(self, key_path='encryption.key'):
        self.key_path = key_path
        self.cipher_suite = self.init_cipher_suite()

    def init_cipher_suite(self):
        key = self.load_or_generate_key()
        return Fernet(key)

    def load_or_generate_key(self):
        if os.path.exists(self.key_path):
            with open(self.key_path, 'rb') as key_file:
                return key_file.read()
        key = Fernet.generate_key()
        with open(self.key_path, 'wb') as key_file:
            key_file.write(key)
        return key

    def encrypt_text(self, text):
        if isinstance(text, str):
            text = text.encode()
        return self.cipher_suite.encrypt(text)

    def decrypt_text(self, encrypted_text):
        return self.cipher_suite.decrypt(encrypted_text).decode()

# PDF to images conversion function
def convert_pdf_to_images(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        images = [convert_from_path(page.path, 200) for page in pdf.pages]
    return images

class SettingsManager:
    def __init__(self, settings_path='settings.json'):
        self.settings_path = settings_path
        self.settings = self.load_settings()

    def load_settings(self):
        if os.path.exists(self.settings_path):
            with open(self.settings_path, 'r') as file:
                return json.load(file)
        return {"custom_patterns": {}}

    def save_settings(self):
        with open(self.settings_path, 'w') as file:
            json.dump(self.settings, file, indent=4)



def read_file_content(file_path):
    try:
        if file_path.endswith('.txt'):
            return read_text_file(file_path)
        elif file_path.endswith('.pdf'):
            return read_pdf_file(file_path)
            try:
                text = ''
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() + ' '
                if not text.strip():
                    images = convert_pdf_to_images(file_path)
                    text = ' '.join([pytesseract.image_to_string(image) for image in images])
                return text
            except Exception as e:
                logging.error(f"Error reading PDF file {file_path}: {e}")
                messagebox.showerror("Error", f"Could not read PDF file: {file_path}")
                return None
        elif file_path.endswith('.docx'):
            doc = docx.Document(file_path)
            text = '\n'.join(paragraph.text for paragraph in doc.paragraphs)
            return text
        elif file_path.endswith('.xlsx'):
            workbook = openpyxl.load_workbook(file_path)
            text = ''
            for sheet in workbook.sheetnames:
                ws = workbook[sheet]
                for row in ws.iter_rows(values_only=True):
                    if row:
                        text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
            return text
        elif file_path.endswith('.pptx'):
            pythoncom.CoInitialize()
            pres = Presentation(file_path)
            text = ''
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
            return text
        elif file_path.endswith('.html'):
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                return soup.get_text()
        elif file_path.endswith('.md'):
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                html = markdown2.markdown(text)
                soup = BeautifulSoup(html, 'html.parser')
                return soup.get_text()
        else:
            logging.warning(f"Unsupported file format for {file_path}")
            messagebox.showerror("Unsupported File", "The selected file format is not supported.")
            return None
    except FileNotFoundError:
        logger.error(f"File not found: {file_path}")
        messagebox.showerror("Error", f"File not found: {file_path}")
        return None
    except PermissionError:
        logger.error(f"Permission denied for file: {file_path}")
        messagebox.showerror("Error", f"Permission denied for file: {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error reading file {file_path}: {e}")
        messagebox.showerror("Error", f"Could not read file: {file_path}")
        return None
    
def read_text_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_pdf_file(file_path):
    text = ''
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + ' '
    if not text.strip():
        images = convert_pdf_to_images(file_path)
        with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
            future_to_image = {executor.submit(pytesseract.image_to_string, image): image for image in images}
            for future in concurrent.futures.as_completed(future_to_image):
                text += future.result() + ' '
    return text

class SettingsManager:
    def __init__(self, settings_path='settings.json'):
        self.settings_path = settings_path

    def load_settings(self):
        if os.path.exists(self.settings_path):
            with open(self.settings_path, 'r') as file:
                return json.load(file)
        return self.settings_manager.load_settings()

    def save_settings(self, settings):
        with open(self.settings_path, 'w') as file:
            json.dump(settings, file, indent=4)

def compile_regex_patterns(custom_patterns=None):
    default_patterns = {
        'Email Address': re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'),
        'Phone Number': re.compile(r'\b(?:\+?1[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b'),
        'Address': re.compile(r'\b\d+\s+[\w\s]+,\s+\w+\b|\b\d+\s+[\w\s]+\b'),
        'Social Security Number': re.compile(r'\b\d{3}-?\d{2}-?\d{4}\b'),
        'Credit Card Number': re.compile(r'\b\d{4}-?\d{4}-?\d{4}-?\d{4}\b'),
        'Passport Number': re.compile(r'\b[A-Z]{1,2}\d{6,9}\b'),
        'Driver\'s License Number': re.compile(r'\b[Ss]\d{7,8}\b'),
    }
    if custom_patterns:
        default_patterns.update(custom_patterns)
    return default_patterns

def detect_pii(file_path, selected_pii, custom_patterns=None):
    content = read_file_content(file_path)
    if content is None:
        logging.info(f"No content read from {file_path}")
        return {}
    
    regex_patterns = compile_regex_patterns(custom_patterns)
    detected_pii = {}
    
    for pattern_name, pattern in regex_patterns.items():
        if pattern_name in selected_pii:
            matches = pattern.findall(content)
            if matches:
                detected_pii[pattern_name] = matches
                logging.info(f"Detected {pattern_name} in {file_path}")
    if not detected_pii:
        logging.info(f"No PII detected in {file_path}")
    return detected_pii

# PII Scanner App class definition with progress_var encapsulation
class PiiScannerApp:
    def __init__(self, master):
        self.master = master
        master.title("PII Detector")

        self.settings_manager = SettingsManager()  # Initialize SettingsManager
        self.settings = self.settings_manager.load_settings()  # Use SettingsManager to load settings
        self.custom_patterns = self.settings.get('custom_patterns', {})  # Initialize custom_patterns from settings

        self.encryption_module = EncryptionModule()

        self.pii_selections = []

        self.progress_var = tk.DoubleVar()  # Progress variable for progress bar

        self.setup_ui()

        global progress_var
        progress_var = tk.DoubleVar()
        self.progress_bar = Progressbar(master, orient="horizontal", length=200, mode="determinate", variable=progress_var)
        self.progress_bar.pack()
    
    def setup_ui(self):
        self.master.title("PII Scanner")
        self.master.geometry("800x600")

        self.style = Style(self.master)
        self.style.configure('TCheckbutton', font=('Arial', 10))

        self.main_frame = tk.Frame(self.master)
        self.main_frame.pack(fill=tk.BOTH, expand=True)

        self.top_frame = tk.Frame(self.main_frame)
        self.top_frame.pack(fill=tk.X)

        self.bottom_frame = tk.Frame(self.main_frame)
        self.bottom_frame.pack(fill=tk.BOTH, expand=True)

        self.file_button = tk.Button(self.top_frame, text="Select Files", command=self.browse_files)
        self.file_button.pack(side=tk.LEFT, padx=10, pady=10)

        self.scan_button = tk.Button(self.top_frame, text="Scan", command=self.scan_files)
        self.scan_button.pack(side=tk.LEFT, padx=10)

        self.save_button = tk.Button(self.top_frame, text="Save Report", command=self.save_report_dialog)
        self.save_button.pack(side=tk.LEFT, padx=10)

        self.settings_button = Button(self.master, text="Settings", command=self.open_settings)
        self.settings_button.pack()

        self.separator = Separator(self.main_frame, orient='horizontal')
        self.separator.pack(fill=tk.X, padx=5, pady=5)

        self.check_frame = tk.Frame(self.bottom_frame)
        self.check_frame.pack(side=tk.LEFT, fill=tk.Y, padx=10)

        self.scrollbar = tk.Scrollbar(self.check_frame)
        self.scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        self.progress_bar = Progressbar(self.master, orient="horizontal", length=200, mode="determinate", variable=self.progress_var)
        self.progress_bar.pack()

        self.pii_listbox = tk.Listbox(self.check_frame, yscrollcommand=self.scrollbar.set, selectmode=tk.MULTIPLE)
        self.pii_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self.scrollbar.config(command=self.pii_listbox.yview)

        self.results_tree = ttk.Treeview(self.bottom_frame, columns=("Type", "Value"), show="headings")
        self.results_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=10, pady=10)
        self.results_tree.heading("Type", text="PII Type")
        self.results_tree.heading("Value", text="PII Value")

    # Initialize the frame for Checkbuttons here
        self.checkbuttons_frame = tk.Frame(self.check_frame)
        self.checkbuttons_frame.pack(side=tk.LEFT, fill=tk.Y, padx=10)

        self.pii_types = [
            'Email Address', 'Phone Number', 'Address',
            'Social Security Number', 'Credit Card Number', 'Passport Number', 'Driver\'s License Number'
        ] + list(self.custom_patterns.keys())
        for pii_type in self.pii_types:
            var = IntVar()
            chk = Checkbutton(self.checkbuttons_frame, text=pii_type, variable=var, font=('Arial', 10))
            chk.pack(anchor=tk.W, side=tk.TOP, expand=True)
            self.pii_selections.append((pii_type, var))

        self.pii_listbox.bind('<<ListboxSelect>>', self.on_select)

    # Added Text widget for displaying results
        self.results_text = tk.Text(self.bottom_frame)
        self.results_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=10, pady=10)

        pass
       
    
    def on_select(self, event):
        selections = event.widget.curselection()
        self.selected_pii = [self.pii_types[i] for i in selections]
        
        

    def open_settings(self):
        settings_window = Toplevel(self.master)
        settings_window.title("Settings")
        settings_window.geometry("400x300")

        self.custom_patterns_listbox = Listbox(settings_window)
        self.custom_patterns_listbox.pack(pady=5)

        self.custom_pii_name_entry = Entry(settings_window)
        self.custom_pii_name_entry.pack(pady=5)

        self.custom_pii_pattern_entry = Entry(settings_window)
        self.custom_pii_pattern_entry.pack(pady=5)

        Button(settings_window, text="Add Pattern", command=self.add_custom_pattern).pack(pady=2)
        Button(settings_window, text="Remove Selected Pattern", command=self.remove_custom_pattern).pack(pady=2)

        # Update the listbox with current custom patterns
        self.update_custom_patterns_listbox()
        
    def remove_custom_pattern(self):
            selected_indices = self.custom_patterns_listbox.curselection()
            if not selected_indices:
                messagebox.showerror("Error", "No pattern selected.")
                return
            selected_index = selected_indices[0]
            pattern_name = self.custom_patterns_listbox.get(selected_index)
            del self.custom_patterns[pattern_name]  # Remove the pattern from the dictionary
            self.settings['custom_patterns'] = self.custom_patterns  # Update the settings dictionary
            self.settings_manager.save_settings(self.settings)  # Save the updated settings to the JSON file
            self.update_custom_patterns_listbox()  # Refresh the listbox to show current custom patterns
            messagebox.showinfo("Success", f"Removed pattern '{pattern_name}' successfully.")
    
    def save_custom_pattern(self):
        pattern_name = self.custom_pii_name_entry.get()
        pattern_regex = self.custom_pii_pattern_entry.get()
        if pattern_name and pattern_regex:
        # Assuming you have a way to update self.settings with the new pattern
            self.custom_patterns[pattern_name] = pattern_regex
            self.settings['custom_patterns'] = self.custom_patterns
            self.settings_manager.save_settings(self.settings)
            messagebox.showinfo("Success", "Custom PII pattern saved.")
            self.settings_window.destroy()
            self.setup_ui() 
            
    def browse_files(self):
        self.file_paths = filedialog.askopenfilenames()
        if not self.file_paths:
            messagebox.showinfo("Selection", "No files selected.")
    
    def scan_files(self):
        selected_pii = [pii_type for pii_type, var in self.pii_selections if var.get()]
        if self.file_paths and selected_pii:
            logging.info("Scan started.")
            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                future_to_file = {executor.submit(detect_pii, file_path, selected_pii, self.custom_patterns): file_path for file_path in self.file_paths}
                for future in concurrent.futures.as_completed(future_to_file):
                    file_path = future_to_file[future]
                    try:
                        result = future.result()
                        self.display_results(result, file_path)
                    except Exception as e:
                        logging.error(f"Error scanning {file_path}: {e}")
                        messagebox.showerror("Scan Error", f"An error occurred while scanning {file_path}.")
            logging.info("Scan completed.")
        else:
            messagebox.showwarning("Warning", "Please select one or more files and at least one PII type.")
            threading.Thread(target=self._scan_files_thread, daemon=True).start()
    
            
       
    def display_results(self, results, file_path):
        self.results_text.insert(tk.END, f"Results for {os.path.basename(file_path)}:\n")
        if not results:
            self.results_text.insert(tk.END, "No PII found.\n\n")
        else:
            for pii_type, piis in results.items():
                parent = self.results_tree.insert("", tk.END, text=pii_type, values=(f"{os.path.basename(file_path)} - {pii_type}",))
            for pii in piis:
                self.results_tree.insert(parent, tk.END, values=(" ", pii))
    
    def save_report_dialog(self):
        format_choice = simpledialog.askstring("Save Report", "Enter format (TXT, PDF, CSV):")
        if format_choice:
            self.save_report(format_choice.upper())
    
    def save_report(self, format_choice):
        content = self.results_text.get(1.0, tk.END)
        # Use the encrypt_text method from the EncryptionModule instance
        encrypted_content = self.encryption_module.encrypt_text(content)
        
        save_path = filedialog.asksaveasfilename(defaultextension=f".{format_choice.lower()}")
        if save_path:
            if format_choice == 'TXT':
                with open(save_path, 'wb') as file:
                    file.write(encrypted_content)
            elif format_choice == 'PDF' or format_choice == 'CSV':
                messagebox.showerror("Error", "Encryption currently supported only for TXT format.")
            messagebox.showinfo("Info", "Report saved successfully.")
            
    def test_pattern(self):
        test_string = simpledialog.askstring("Test Custom Pattern", "Enter text to test pattern:")
        if test_string:
            try:
                pattern = re.compile(self.custom_pii_pattern_entry.get())
                matches = pattern.findall(test_string)
                if matches:
                    messagebox.showinfo("Pattern Test Result", f"Matches found: {', '.join(matches)}")
                else:
                    messagebox.showinfo("Pattern Test Result", "No matches found.")
            except re.error as e:
                messagebox.showerror("Pattern Test Error", f"Invalid regex pattern: {e}")
                
    def add_custom_pattern(self):
        pattern_name = self.custom_pii_name_entry.get()
        pattern_regex = self.custom_pii_pattern_entry.get()
        if pattern_name and pattern_regex:
            try:
                re.compile(pattern_regex)  # Validate regex
                self.custom_patterns[pattern_name] = pattern_regex
                self.settings['custom_patterns'] = self.custom_patterns
                self.settings_manager.save_settings()  # Corrected usage
                messagebox.showinfo("Success", "Pattern added successfully.")
            except re.error as e:
                messagebox.showerror("Error", f"Invalid regex pattern: {e}")
                
    def update_custom_patterns_listbox(self):
        # Assuming self.custom_patterns_listbox is correctly initialized
        self.custom_patterns_listbox.delete(0, tk.END)
        for pattern_name in self.custom_patterns.keys():
            self.custom_patterns_listbox.insert(tk.END, pattern_name)
                

@app.route('/api/scan', methods=['POST'])
@auth.login_required
def api_scan():
    data = flask.request.json
    file_path = data['file_path']
    selected_pii = data['selected_pii']
    # Assuming detect_pii function exists and is properly defined
    result = detect_pii(file_path, selected_pii)
    return flask.jsonify(result)

if __name__ == "__main__":
    # Setup logging configuration (optional but recommended for better logging behavior)
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

    api_thread = Thread(target=lambda: app.run(port=5000, debug=True, use_reloader=False))
    api_thread.start()
    root = tk.Tk()
    app = PiiScannerApp(root)
    root.mainloop()